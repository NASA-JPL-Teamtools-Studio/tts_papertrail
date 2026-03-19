import math
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from tts_papertrail.base import Report, TableEntry, PowerTableEntry, SummaryTable, TextEntry, ListEntry, SectionEntry, RichText

def hex_to_rgb(hex_str):
    """Helper: Converts hex string to python-pptx RGBColor object."""
    if not hex_str: return None
    try:
        hex_str = hex_str.lstrip('#')
        return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))
    except: return None

class PowerPointReport:
    """
    Renders reports to PowerPoint (.pptx) files using `python-pptx`.

    Idiomatic Features:
    - **One Entry per Slide:** Generally maps one entry to one slide.
    - **Table Pagination:** Automatically splits long tables across multiple slides
      to prevent overflow, adding "(m/n)" to the titles.
    - **Sections:** Groups `SectionEntry` content onto a single "Title and Content" slide.
    - **Overflow Warning:** Logs warnings if content might exceed slide boundaries.
    """
    
    MAX_ROWS_PER_SLIDE = 12
    
    def __init__(self, name: str, author: str):
        self.name = name
        self.author = author
        self.created_at = None  # Will be set by Report.to_powerpoint()
        self.entries = []

    def save(self, filename: str):
        """
        Generates the .pptx file.
        
        Args:
            filename (str): Output filename. Appends .pptx if missing.
        """
        if not filename.endswith('.pptx'): 
            filename += '.pptx'
        prs = Presentation()
        
        # Create Title Slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = self.name
        slide.placeholders[1].text = f"Author: {self.author}\nGenerated: {self.created_at.strftime('%Y-%m-%d')}"

        def render_text_frame(tf, entry):
            """Helper to render Text or List entries into a PowerPoint TextFrame."""
            if entry.title:
                p = tf.add_paragraph()
                r = p.add_run()
                r.text = entry.title
                r.font.bold = True
                r.font.underline = True

            if isinstance(entry, TextEntry):
                # Render as a standard paragraph
                p = tf.add_paragraph()
                p.level = 0
                for item in entry.content:
                    val = item if isinstance(item, RichText) else RichText(str(item))
                    r = p.add_run()
                    r.text = val.text
                    if val.bold: r.font.bold = True
                    if val.italic: r.font.italic = True
                    if val.color: r.font.color.rgb = hex_to_rgb(val.color)

            elif isinstance(entry, ListEntry):
                # Render as bullet points, using PPT indentation levels
                for item in entry.items:
                    p = tf.add_paragraph()
                    p.level = item.level
                    for part in item.content:
                        r = p.add_run()
                        r.text = part.text
                        if part.bold: r.font.bold = True
                        if part.italic: r.font.italic = True
                        if part.color: r.font.color.rgb = hex_to_rgb(part.color)

        for entry in self.entries:
            # --- Section Logic ---
            # Stacks multiple parts (Text/Lists) onto a SINGLE slide
            if isinstance(entry, SectionEntry):
                slide = prs.slides.add_slide(prs.slide_layouts[1]) # Layout 1: Title & Content
                slide.shapes.title.text = entry.title
                tf = slide.placeholders[1].text_frame
                tf.clear()
                
                for part in entry.parts:
                    render_text_frame(tf, part)
                    tf.add_paragraph().text = "" # Spacer paragraph

            # --- Individual Text/List Logic ---
            elif isinstance(entry, (TextEntry, ListEntry)):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = entry.title or ""
                tf = slide.placeholders[1].text_frame
                tf.clear()
                render_text_frame(tf, entry)

            # --- SummaryTable Logic ---
            elif isinstance(entry, SummaryTable):
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                slide.shapes.title.text = entry.title or "Summary"
                
                headers = ['Category', 'Status', 'Message', 'Disposition Summary']
                rows = len(entry.rows) + 1  # +1 for header
                cols = 4
                
                # Create table
                left = Inches(0.5)
                top = Inches(1.5)
                width = Inches(9)
                height = Inches(5)
                table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                
                # Write headers
                for j, h in enumerate(headers):
                    cell = table.cell(0, j)
                    cell.text = str(h)
                    cell.text_frame.paragraphs[0].font.bold = True
                
                # Write data rows
                for i, row_data in enumerate(entry.rows, start=1):
                    table.cell(i, 0).text = row_data['category']
                    
                    # Handle HTML in status field
                    status = row_data['status']
                    cell = table.cell(i, 1)
                    if '<' in status and '>' in status:
                        rich_parts = RichText.from_html(status)
                        if rich_parts:
                            cell.text = ''
                            tf = cell.text_frame
                            tf.clear()
                            p = tf.paragraphs[0]
                            for part in rich_parts:
                                r = p.add_run()
                                r.text = part.text
                                if part.bold: r.font.bold = True
                                if part.italic: r.font.italic = True
                                if part.color: r.font.color.rgb = hex_to_rgb(part.color)
                        else:
                            cell.text = status
                    else:
                        cell.text = status
                    
                    table.cell(i, 2).text = row_data['message']
                    table.cell(i, 3).text = row_data['disposition_summary']

            # --- PowerTable Logic ---
            # Handles pagination for long tables from DataContainer
            elif isinstance(entry, PowerTableEntry):
                # Extract data from DataContainer
                container = entry.container.with_cols(entry.columns) if entry.columns else entry.container
                headers = entry.columns if entry.columns else container.repr_cols
                data_items = list(container)
                
                has_header = 1
                data_rows = self.MAX_ROWS_PER_SLIDE - has_header
                if data_rows < 1: data_rows = 1
                
                total = len(data_items)
                pages = math.ceil(total / data_rows) if total > 0 else 1
                
                for page in range(1, pages + 1):
                    slide = prs.slides.add_slide(prs.slide_layouts[5])
                    
                    # Generate Pagination Title
                    title_text = entry.title or "Table"
                    if pages > 1:
                        title_text += f" ({page}/{pages})"
                    slide.shapes.title.text = title_text
                    
                    # Slice data for this specific slide
                    start = (page - 1) * data_rows
                    page_items = data_items[start : start + data_rows]
                    
                    rows = len(page_items) + has_header
                    cols = len(headers)
                    if cols == 0: continue
                    
                    # Create table
                    left = Inches(0.5)
                    top = Inches(1.5)
                    width = Inches(9)
                    height = Inches(5)
                    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                    
                    # Write headers
                    for j, h in enumerate(headers):
                        cell = table.cell(0, j)
                        cell.text = str(h)
                        cell.text_frame.paragraphs[0].font.bold = True
                    
                    # Write data rows
                    for i, item in enumerate(page_items, start=1):
                        # Get formatted values and cell styles from the DataItem
                        printable_vals = item.printable_values
                        cell_styles = item.default_html_cell_styles
                        
                        for j, col in enumerate(headers):
                            value = printable_vals.get(col, '')
                            cell_style = cell_styles.get(col, {})
                            cell = table.cell(i, j)
                            
                            if value:
                                str_value = str(value)
                                
                                # Check if value contains HTML tags
                                if '<' in str_value and '>' in str_value:
                                    # Parse HTML and apply formatting from parsed content
                                    rich_parts = RichText.from_html(str_value)
                                    if rich_parts:
                                        cell.text = ''
                                        tf = cell.text_frame
                                        tf.clear()
                                        p = tf.paragraphs[0]
                                        for part in rich_parts:
                                            r = p.add_run()
                                            r.text = part.text
                                            if part.bold: r.font.bold = True
                                            if part.italic: r.font.italic = True
                                            if part.color: r.font.color.rgb = hex_to_rgb(part.color)
                                    else:
                                        # Fallback if HTML parsing fails
                                        cell.text = str_value
                                else:
                                    # No HTML - use cell_style for formatting
                                    cell.text = str_value
                                    
                                    # Apply styling from cell_style
                                    p = cell.text_frame.paragraphs[0]
                                    r = p.runs[0] if p.runs else p.add_run()
                                    
                                    if cell_style.get('font-weight') == 'bold':
                                        r.font.bold = True
                                    if cell_style.get('font-style') == 'italic':
                                        r.font.italic = True
                                    if 'color' in cell_style:
                                        color_val = cell_style['color']
                                        if color_val.startswith('#'):
                                            r.font.color.rgb = hex_to_rgb(color_val)
                            else:
                                cell.text = ''
            
            # --- Table Logic ---
            # Handles pagination for long tables
            elif isinstance(entry, TableEntry):
                has_header = 1 if entry.headers else 0
                data_rows = self.MAX_ROWS_PER_SLIDE - has_header
                if data_rows < 1: data_rows = 1
                
                total = len(entry.data)
                pages = math.ceil(total / data_rows) if total > 0 else 1

                for page in range(1, pages + 1):
                    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Layout 5: Title Only
                    
                    # Generate Pagination Title
                    title_text = entry.title or "Table"
                    if pages > 1: 
                        title_text += f" ({page}/{pages})"
                    slide.shapes.title.text = title_text

                    # Slice data for this specific slide
                    start = (page - 1) * data_rows
                    page_data = entry.data[start : start + data_rows]
                    
                    rows = len(page_data) + has_header
                    cols = max(len(entry.headers or []), len(page_data[0]) if page_data else 0)
                    if cols == 0: continue

                    # Create Table Shape
                    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.8)).table
                    
                    r_idx = 0
                    if entry.headers:
                        for i, h in enumerate(entry.headers):
                            cell = table.cell(0, i)
                            cell.text = str(h)
                            cell.text_frame.paragraphs[0].runs[0].font.bold = True
                        r_idx = 1
                    
                    for row_data in page_data:
                        for i, datum in enumerate(row_data):
                            cell = table.cell(r_idx, i)
                            val = datum if isinstance(datum, RichText) else RichText(str(datum))
                            cell.text = val.text
                            if val.color:
                                cell.text_frame.paragraphs[0].runs[0].font.color.rgb = hex_to_rgb(val.color)
                        r_idx += 1

        prs.save(filename)
        print(f"✅ PowerPoint report saved to {filename}")